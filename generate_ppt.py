"""
CEXP Performance Dashboard — PowerPoint Generator
=================================================
INSTRUCTIONS:
  1. Fill in your real data in the DATA section below.
  2. Run:  python3 generate_ppt.py
  3. Open: cexp_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import copy

# ============================================================
# DATA SECTION — Replace with your actual CEXP data
# ============================================================
MONTHS = ["Oct-2025","Nov-2025", "Dec-2025", "Jan-2026", "Feb-2026", "Mar-2026"]

DATA = {
    "Profile Summary": {
        "type": "frontend",
        "label": "Profile Summary",   # <-- Replace with actual name
        "hits":     [53892,31585, 29271, 39966, 38245, 33165],  # <-- Replace with hit counts
        "p90":      [3138, 3162, 3031, 2908, 3087, 3377],  # <-- P90 response times (ms)
        "p95":      [14466, 15028, 11758, 10652, 13340, 14140],  # <-- P95 response times (ms)
        "p99":      [4590, 4630, 4370, 4092, 4324, 4688],  # <-- P99 response times (ms)
        "avg":      [1857, 1862, 1703, 1633, 1777, 2035],  # <-- Average response times (ms)
        "errors500": None,                # Not applicable for frontend
    },
    "Alerts": {
        "type": "frontend",
        "label": "Alerts",
        "hits":     [10841,8502, 7936, 12133, 10165, 9572],
        "p90":      [5001, 4384, 4581, 4113, 4431, 3813],
        "p95":      [21389, 16477, 19638, 15587, 17412, 14922],
        "p99":      [5890, 5936, 6614, 5607, 6021, 5215],
        "avg":      [3041, 2511, 2602, 2406, 2497, 2282],
        "errors500": None,
    },
    "ClientPreferenceSet": {
        "type": "backend",
        "label": "ClientPreferenceSet",
        "hits":     [105805, 66222, 82247, 94415, 88818, 87229],
        "p90":      [786, 831, 894, 805, 773, 776],
        "p95":      [1986, 2922, 1966, 1856, 1885, 1921],
        "p99":      [1031, 1183, 1138, 1039, 1015, 1023],
        "avg":      [47, 42, 1, 1, 1, 27],
        "errors500": [21, 7, 4, 9, 0, 1],
    },
    "MaintainClientPreference": {
        "type": "backend",
        "label": "Maintain Client Preference",
        "hits":     [17324, 10519, 11725, 13110, 12002, 11860],  # <-- Replace with hit counts
        "p90":      [893, 907, 832, 809, 805, 816],  # <-- P90 response times (ms)
        "p95":      [2177, 4912, 1902, 1911, 1864, 1880],  # <-- P95 response times (ms)
        "p99":      [1089, 1221, 1022, 1005, 999, 998],  # <-- P99 response times (ms)
        "avg":      [180, 249, 83, 142, 120, 246],  # <-- Average response times (ms)
        "errors500": [23, 9, 9, 18, 14, 10], # <-- 500 error counts
    },
}
# ============================================================
# END DATA SECTION
# ============================================================

OUTPUT_FILE = "cexp_presentation.pptx"

# ── Colour palette ───────────────────────────────────────────
C = {
    "bg":       RGBColor(0x0F, 0x17, 0x2A),
    "surface":  RGBColor(0x1E, 0x29, 0x3B),
    "surface2": RGBColor(0x33, 0x41, 0x55),
    "accent1":  RGBColor(0x38, 0xBD, 0xF8),   # sky blue   — P90 / FE1
    "accent2":  RGBColor(0xA7, 0x8B, 0xFA),   # violet     — P95 / FE2
    "accent3":  RGBColor(0x34, 0xD3, 0x99),   # emerald    — Avg / BE1
    "accent4":  RGBColor(0xFB, 0x92, 0x3C),   # orange     — Hits
    "pink":     RGBColor(0xF4, 0x72, 0xB6),   # pink       — P99
    "red":      RGBColor(0xF8, 0x71, 0x71),   # red        — Errors
    "yellow":   RGBColor(0xFB, 0xBF, 0x24),   # amber      — Warn
    "white":    RGBColor(0xF1, 0xF5, 0xF9),
    "muted":    RGBColor(0x94, 0xA3, 0xB8),
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────
def safe_avg(arr):
    vals = [v for v in arr if v and v > 0]
    return round(sum(vals) / len(vals)) if vals else 0

def safe_min(arr):
    vals = [v for v in arr if v and v > 0]
    return min(vals) if vals else 0

def safe_max(arr):
    vals = [v for v in arr if v and v > 0]
    return max(vals) if vals else 0

def fmt_num(v):
    if not v: return "—"
    if v >= 1_000_000: return f"{v/1_000_000:.1f}M"
    if v >= 1_000: return f"{v/1_000:.1f}K"
    return str(int(v))

def fmt_ms(v):
    return f"{v} ms" if v else "—"

def calc_trend(arr):
    vals = [v for v in arr if v and v > 0]
    if len(vals) < 2:
        return "stable"
    half = max(1, len(vals) // 2)
    first = sum(vals[:half]) / half
    last  = sum(vals[-half:]) / half
    pct   = (last - first) / first if first else 0
    if pct < -0.05: return "improving"
    if pct >  0.05: return "degrading"
    return "stable"

def trend_label(t):
    return {"improving": "↓ Improving", "degrading": "↑ Degrading", "stable": "→ Stable"}[t]

def trend_color(t):
    return {"improving": C["accent3"], "degrading": C["red"], "stable": C["muted"]}[t]


# ── Low-level drawing helpers ─────────────────────────────────
def fill_slide_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill=None, line=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=Pt(11), bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or C["white"]
    return txb

def add_label(slide, text, left, top, width):
    """Small all-caps section label."""
    add_text(slide, text.upper(), left, top, width, Inches(0.2),
             font_size=Pt(8), bold=True, color=C["muted"])

def add_tag(slide, text, left, top, color=None):
    """Small coloured tag pill."""
    color = color or C["accent1"]
    add_text(slide, text, left, top, Inches(2), Inches(0.25),
             font_size=Pt(8), bold=True, color=color)

def add_accent_line(slide, left, top, width=Inches(0.6)):
    bar = add_rect(slide, left, top, width, Inches(0.04), fill=C["accent1"])
    return bar

def kpi_card(slide, left, top, w, h, value, label, sub=None, val_color=None):
    add_rect(slide, left, top, w, h, fill=C["surface"], line=C["surface2"])
    add_text(slide, value, left, top + Inches(0.12), w, Inches(0.35),
             font_size=Pt(20), bold=True, color=val_color or C["accent1"], align=PP_ALIGN.CENTER)
    add_text(slide, label, left, top + Inches(0.46), w, Inches(0.2),
             font_size=Pt(8), bold=True, color=C["muted"], align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, sub, left, top + Inches(0.64), w, Inches(0.18),
                 font_size=Pt(7), color=C["muted"], align=PP_ALIGN.CENTER)

def callout_box(slide, left, top, w, h, title, body, style="good"):
    style_colors = {"good": C["accent3"], "warn": C["yellow"], "bad": C["red"]}
    col = style_colors.get(style, C["accent3"])
    add_rect(slide, left, top, w, h, fill=C["surface"], line=C["surface2"])
    # left border accent
    add_rect(slide, left, top, Inches(0.04), h, fill=col)
    add_text(slide, title, left + Inches(0.1), top + Inches(0.06), w - Inches(0.15), Inches(0.18),
             font_size=Pt(8), bold=True, color=col)
    add_text(slide, body, left + Inches(0.1), top + Inches(0.26), w - Inches(0.15), h - Inches(0.32),
             font_size=Pt(9), color=C["white"], wrap=True)

def add_line_chart(slide, left, top, width, height, series_data, months, colors):
    """series_data = list of (name, values_list)"""
    cd = ChartData()
    cd.categories = months
    for name, vals in series_data:
        cd.add_series(name, vals)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, left, top, width, height, cd
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    # Style series colours
    for i, (_, _) in enumerate(series_data):
        s = chart.series[i]
        s.format.line.color.rgb = colors[i % len(colors)]
        s.format.line.width = Pt(2)
        s.marker.format.fill.solid()
        s.marker.format.fill.fore_color.rgb = colors[i % len(colors)]
    # Style plot area
    _style_chart_bg(chart)
    _style_axes(chart)
    return chart

def add_bar_chart(slide, left, top, width, height, series_data, months, colors, stacked=False):
    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    cd = ChartData()
    cd.categories = months
    for name, vals in series_data:
        cd.add_series(name, vals)
    chart = slide.shapes.add_chart(chart_type, left, top, width, height, cd).chart
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    for i, _ in enumerate(series_data):
        s = chart.series[i]
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = colors[i % len(colors)]
    _style_chart_bg(chart)
    _style_axes(chart)
    return chart

def _style_chart_bg(chart):
    from lxml import etree
    # Set chart area and plot area fill to dark via XML
    try:
        ca = chart.chart_area
        ca.format.fill.solid()
        ca.format.fill.fore_color.rgb = C["surface"]
    except Exception:
        pass
    try:
        # plot_area via XML
        spPr = chart._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/chart}plotArea')
        if spPr is not None:
            nsmap = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            spPr_el = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
            if spPr_el is None:
                spPr_el = etree.SubElement(spPr, f'{{{nsmap}}}spPr')
            solidFill = etree.SubElement(spPr_el, f'{{{nsmap}}}solidFill')
            srgb = etree.SubElement(solidFill, f'{{{nsmap}}}srgbClr')
            srgb.set('val', '1E293B')
    except Exception:
        pass

def _style_axes(chart):
    try:
        va = chart.value_axis
        va.format.line.color.rgb = C["surface2"]
        va.tick_labels.font.color.rgb = C["muted"]
        va.tick_labels.font.size = Pt(8)
    except Exception:
        pass
    try:
        ca = chart.category_axis
        ca.format.line.color.rgb = C["surface2"]
        ca.tick_labels.font.color.rgb = C["muted"]
        ca.tick_labels.font.size = Pt(8)
    except Exception:
        pass


# ── Slide builders ────────────────────────────────────────────

def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill_slide_bg(slide, C["bg"])

    # Decorative gradient rect on left
    add_rect(slide, 0, 0, Inches(4.5), SLIDE_H, fill=RGBColor(0x1A, 0x32, 0x50))

    # Tag
    add_text(slide, "CEXP PERFORMANCE REPORT", Inches(0.5), Inches(1.6), Inches(6), Inches(0.3),
             font_size=Pt(9), bold=True, color=C["accent1"])

    # Title
    add_text(slide, "Product Performance Review", Inches(0.5), Inches(2.1), Inches(8), Inches(0.65),
             font_size=Pt(38), bold=True, color=C["white"])

    # Accent line
    add_accent_line(slide, Inches(0.5), Inches(2.9), Inches(0.8))

    # Subtitle
    add_text(slide, "Top 3 Business Functions — Last 6 Months",
             Inches(0.5), Inches(3.1), Inches(8), Inches(0.35),
             font_size=Pt(16), color=C["muted"])

    # Meta grid
    meta = [
        ("PERIOD", "Oct 2025 – Mar 2026"),
        ("FUNCTIONS", "3 Key Functions"),
        ("DATA SOURCE", "CEXP Tool"),
        ("PREPARED", "April 2026"),
    ]
    for i, (lbl, val) in enumerate(meta):
        x = Inches(0.5 + i * 2.9)
        add_text(slide, lbl, x, Inches(4.2), Inches(2.5), Inches(0.2),
                 font_size=Pt(8), color=C["muted"])
        add_text(slide, val, x, Inches(4.45), Inches(2.5), Inches(0.25),
                 font_size=Pt(12), bold=True, color=C["white"])

    # Decorative dots
    for i in range(5):
        add_rect(slide, Inches(10.5 + i*0.35), Inches(6.8), Inches(0.18), Inches(0.18),
                 fill=C["surface2"])


def build_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide, C["bg"])

    add_text(slide, "OVERVIEW", Inches(0.5), Inches(0.3), Inches(4), Inches(0.22),
             font_size=Pt(8), bold=True, color=C["accent1"])
    add_text(slide, "Executive Summary", Inches(0.5), Inches(0.55), Inches(10), Inches(0.5),
             font_size=Pt(26), bold=True, color=C["white"])
    add_accent_line(slide, Inches(0.5), Inches(1.1))

    n = len(DATA)
    card_w = Inches(11.8 / n - 0.2)
    card_gap = Inches(11.8 / n)
    card_h = Inches(5.5)
    card_tops = Inches(1.3)

    for i, (key, d) in enumerate(DATA.items()):
        cx = Inches(0.5 + i * float(card_gap))
        add_rect(slide, cx, card_tops, card_w, card_h, fill=C["surface"], line=C["surface2"])

        # Type badge color
        is_fe = d["type"] == "frontend"
        badge_col = C["accent1"] if is_fe else C["accent2"]
        badge_txt = "FRONTEND · URL" if is_fe else "BACKEND · API"
        add_text(slide, badge_txt, cx + Inches(0.15), card_tops + Inches(0.15),
                 Inches(2.5), Inches(0.22), font_size=Pt(8), bold=True, color=badge_col)

        add_text(slide, d["label"], cx + Inches(0.15), card_tops + Inches(0.45),
                 card_w - Inches(0.3), Inches(0.35), font_size=Pt(14), bold=True, color=C["white"])

        trend = calc_trend(d["avg"])
        add_text(slide, trend_label(trend), cx + Inches(0.15), card_tops + Inches(0.88),
                 Inches(2), Inches(0.22), font_size=Pt(10), bold=True, color=trend_color(trend))

        rows = [
            ("Avg Response", fmt_ms(safe_avg(d["avg"]))),
            ("P90",          fmt_ms(safe_avg(d["p90"]))),
            ("P95",          fmt_ms(safe_avg(d["p95"]))),
            ("P99",          fmt_ms(safe_avg(d["p99"]))),
            ("Total Hits",   fmt_num(sum(d["hits"]))),
        ]
        if d["errors500"]:
            rows.append(("500 Errors", fmt_num(sum(d["errors500"]))))

        for j, (lbl, val) in enumerate(rows):
            ry = card_tops + Inches(1.25 + j * 0.55)
            add_rect(slide, cx + Inches(0.1), ry, card_w - Inches(0.2), Inches(0.5),
                     fill=C["surface"], line=C["surface2"])
            add_text(slide, lbl, cx + Inches(0.2), ry + Inches(0.08),
                     Inches(1.8), Inches(0.3), font_size=Pt(9), color=C["muted"])
            add_text(slide, val, cx + Inches(2.0), ry + Inches(0.08),
                     Inches(1.7), Inches(0.3), font_size=Pt(10), bold=True,
                     color=C["white"], align=PP_ALIGN.RIGHT)


def build_function_slide(prs, key, slide_tag, is_backend=False):
    d = DATA[key]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide, C["bg"])

    badge = "BACKEND PERFORMANCE" if is_backend else "FRONTEND PERFORMANCE"
    badge_col = C["accent2"] if is_backend else C["accent1"]
    pill_txt = "BACKEND · API" if is_backend else "FRONTEND · URL"

    add_text(slide, badge, Inches(0.5), Inches(0.25), Inches(5), Inches(0.22),
             font_size=Pt(8), bold=True, color=badge_col)

    add_text(slide, d["label"], Inches(0.5), Inches(0.5), Inches(8.5), Inches(0.5),
             font_size=Pt(24), bold=True, color=C["white"])

    add_text(slide, pill_txt, Inches(9.2), Inches(0.55), Inches(1.8), Inches(0.28),
             font_size=Pt(8), bold=True, color=badge_col)

    trend = calc_trend(d["avg"])
    add_text(slide, trend_label(trend), Inches(11.1), Inches(0.55), Inches(1.8), Inches(0.28),
             font_size=Pt(9), bold=True, color=trend_color(trend))

    add_accent_line(slide, Inches(0.5), Inches(1.05))

    # ── KPI row ──
    kpi_data = [
        (fmt_num(sum(d["hits"])),    "6-Month Hits",    "Total Traffic"),
        (fmt_ms(safe_avg(d["avg"])), "Avg Response",    "6-month average"),
        (fmt_ms(safe_avg(d["p90"])), "P90 Latency",     "6-month average"),
        (fmt_ms(safe_avg(d["p99"])), "P99 Latency",     "6-month average"),
    ]
    kpi_w = Inches(2.9)
    kpi_h = Inches(0.85)
    for i, (val, lbl, sub) in enumerate(kpi_data):
        kpi_card(slide, Inches(0.5 + i * 3.1), Inches(1.15), kpi_w, kpi_h, val, lbl, sub)

    # ── Latency line chart ──
    chart_left  = Inches(0.5)
    chart_top   = Inches(2.15)
    chart_width = Inches(8.2)
    chart_height = Inches(4.0) if not is_backend else Inches(3.5)

    add_label(slide, "Response Time Trends (ms)", chart_left, chart_top - Inches(0.22), chart_width)
    add_rect(slide, chart_left, chart_top, chart_width, chart_height,
             fill=C["surface"], line=C["surface2"])

    add_line_chart(slide, chart_left + Inches(0.05), chart_top + Inches(0.05),
                   chart_width - Inches(0.1), chart_height - Inches(0.1),
                   [("P90", d["p90"]), ("P95", d["p95"]), ("P99", d["p99"]), ("Avg", d["avg"])],
                   MONTHS,
                   [C["accent1"], C["accent2"], C["pink"], C["accent3"]])

    # ── Right column ──
    rc_left  = Inches(8.85)
    rc_width = Inches(4.0)

    # Hits bar chart
    hits_h = Inches(1.9) if not is_backend else Inches(1.5)
    add_label(slide, "Monthly Hits (Traffic)", rc_left, chart_top - Inches(0.22), rc_width)
    add_rect(slide, rc_left, chart_top, rc_width, hits_h, fill=C["surface"], line=C["surface2"])
    add_bar_chart(slide, rc_left + Inches(0.05), chart_top + Inches(0.05),
                  rc_width - Inches(0.1), hits_h - Inches(0.1),
                  [("Hits", d["hits"])], MONTHS, [C["accent4"]])

    # Errors bar chart (backend only)
    if is_backend and d["errors500"]:
        err_top = chart_top + hits_h + Inches(0.15)
        err_h   = Inches(1.5)
        add_label(slide, "500 Errors per Month", rc_left, err_top - Inches(0.22), rc_width)
        add_rect(slide, rc_left, err_top, rc_width, err_h, fill=C["surface"], line=C["surface2"])
        add_bar_chart(slide, rc_left + Inches(0.05), err_top + Inches(0.05),
                      rc_width - Inches(0.1), err_h - Inches(0.1),
                      [("500 Errors", d["errors500"])], MONTHS, [C["red"]])
        callout_top = err_top + err_h + Inches(0.15)
    else:
        callout_top = chart_top + hits_h + Inches(0.15)

    # Callouts
    best_idx  = d["avg"].index(safe_min(d["avg"])) if safe_min(d["avg"]) else 0
    worst_idx = d["avg"].index(safe_max(d["avg"])) if safe_max(d["avg"]) else 0
    cb_h = Inches(0.7)

    callout_box(slide, rc_left, callout_top, rc_width, cb_h,
                "BEST MONTH",
                f"{MONTHS[best_idx]} — Avg {fmt_ms(d['avg'][best_idx])}, P99 {fmt_ms(d['p99'][best_idx])}",
                "good")

    callout_box(slide, rc_left, callout_top + cb_h + Inches(0.1), rc_width, cb_h,
                "HIGHEST LATENCY MONTH",
                f"{MONTHS[worst_idx]} — Avg {fmt_ms(d['avg'][worst_idx])}, P99 {fmt_ms(d['p99'][worst_idx])}",
                "warn")


def build_comparison_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide, C["bg"])

    add_text(slide, "COMPARISON", Inches(0.5), Inches(0.25), Inches(5), Inches(0.22),
             font_size=Pt(8), bold=True, color=C["accent1"])
    add_text(slide, "All Functions — Side by Side", Inches(0.5), Inches(0.5),
             Inches(10), Inches(0.5), font_size=Pt(24), bold=True, color=C["white"])
    add_accent_line(slide, Inches(0.5), Inches(1.05))

    keys   = list(DATA.keys())
    colors = [C["accent1"], C["accent2"], C["accent3"], C["accent4"]]
    chart_top = Inches(1.2)
    chart_h   = Inches(3.4)

    # Avg latency overlay
    add_label(slide, "Avg Response Time Trend (ms) — All Functions",
              Inches(0.5), chart_top - Inches(0.22), Inches(6))
    add_rect(slide, Inches(0.5), chart_top, Inches(6.0), chart_h,
             fill=C["surface"], line=C["surface2"])
    add_line_chart(slide, Inches(0.55), chart_top + Inches(0.05),
                   Inches(5.9), chart_h - Inches(0.1),
                   [(DATA[k]["label"], DATA[k]["avg"]) for k in keys],
                   MONTHS, colors)

    # Hits comparison
    add_label(slide, "Traffic Volume (Hits) — All Functions",
              Inches(6.8), chart_top - Inches(0.22), Inches(6))
    add_rect(slide, Inches(6.8), chart_top, Inches(6.0), chart_h,
             fill=C["surface"], line=C["surface2"])
    add_bar_chart(slide, Inches(6.85), chart_top + Inches(0.05),
                  Inches(5.9), chart_h - Inches(0.1),
                  [(DATA[k]["label"], DATA[k]["hits"]) for k in keys],
                  MONTHS, colors)

    # Bottom summary cards
    n = len(keys)
    card_w = Inches(11.8 / n - 0.2)
    card_gap = Inches(11.8 / n)
    card_h = Inches(1.5)
    card_top = chart_top + chart_h + Inches(0.2)

    for i, k in enumerate(keys):
        d = DATA[k]
        cx = Inches(0.5 + i * float(card_gap))
        add_rect(slide, cx, card_top, card_w, card_h, fill=C["surface"], line=C["surface2"])
        # top accent bar
        add_rect(slide, cx, card_top, card_w, Inches(0.05), fill=colors[i])
        add_text(slide, d["label"], cx + Inches(0.15), card_top + Inches(0.12),
                 card_w - Inches(0.3), Inches(0.28), font_size=Pt(11), bold=True, color=C["white"])
        trend = calc_trend(d["avg"])
        add_text(slide, trend_label(trend), cx + Inches(0.15), card_top + Inches(0.45),
                 Inches(1.8), Inches(0.22), font_size=Pt(9), bold=True, color=trend_color(trend))
        add_text(slide, f"Avg: {fmt_ms(safe_avg(d['avg']))}   P99: {fmt_ms(safe_avg(d['p99']))}",
                 cx + Inches(0.15), card_top + Inches(0.75), card_w - Inches(0.3), Inches(0.22),
                 font_size=Pt(9), color=C["muted"])
        add_text(slide, f"Hits: {fmt_num(sum(d['hits']))}",
                 cx + Inches(0.15), card_top + Inches(1.0), card_w - Inches(0.3), Inches(0.22),
                 font_size=Pt(9), color=C["muted"])


def build_takeaways_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide, C["bg"])

    add_text(slide, "CONCLUSION", Inches(0.5), Inches(0.25), Inches(5), Inches(0.22),
             font_size=Pt(8), bold=True, color=C["accent1"])
    add_text(slide, "Key Takeaways & Next Steps", Inches(0.5), Inches(0.5),
             Inches(10), Inches(0.5), font_size=Pt(24), bold=True, color=C["white"])
    add_accent_line(slide, Inches(0.5), Inches(1.05))

    strengths = []
    attention = []

    for k, d in DATA.items():
        if all(v == 0 for v in d["avg"]):
            continue
        trend = calc_trend(d["avg"])
        if trend == "improving":
            strengths.append(f"{d['label']}: Response times improving over the 6-month period.")
        elif trend == "stable":
            strengths.append(f"{d['label']}: Stable performance — no significant degradation.")
        else:
            attention.append(f"{d['label']}: Upward latency trend detected — investigate root cause.")

        if d["errors500"]:
            total_err = sum(d["errors500"])
            if total_err == 0:
                strengths.append(f"{d['label']}: Zero 500 errors recorded over 6 months.")
            else:
                attention.append(f"{d['label']}: {fmt_num(total_err)} total 500 errors — review monthly patterns.")

        if safe_avg(d["p99"]) > safe_avg(d["p90"]) * 2:
            attention.append(f"{d['label']}: High P99 vs P90 gap — tail latency spikes may impact users.")

    if not strengths:
        strengths = ["Fill in your data to see auto-generated insights."]
    if not attention:
        attention = ["Fill in your data to see auto-generated insights."]

    # Strengths column
    add_text(slide, "✓  Strengths", Inches(0.5), Inches(1.25), Inches(5.8), Inches(0.3),
             font_size=Pt(13), bold=True, color=C["accent3"])
    for i, s in enumerate(strengths[:4]):
        iy = Inches(1.65 + i * 0.82)
        add_rect(slide, Inches(0.5), iy, Inches(5.8), Inches(0.72),
                 fill=C["surface"], line=C["surface2"])
        add_rect(slide, Inches(0.5), iy, Inches(0.05), Inches(0.72), fill=C["accent3"])
        add_text(slide, s, Inches(0.7), iy + Inches(0.1), Inches(5.5), Inches(0.55),
                 font_size=Pt(9.5), color=C["white"], wrap=True)

    # Areas to watch column
    add_text(slide, "⚠  Areas to Watch", Inches(6.8), Inches(1.25), Inches(5.8), Inches(0.3),
             font_size=Pt(13), bold=True, color=C["yellow"])
    for i, s in enumerate(attention[:4]):
        iy = Inches(1.65 + i * 0.82)
        add_rect(slide, Inches(6.8), iy, Inches(5.8), Inches(0.72),
                 fill=C["surface"], line=C["surface2"])
        add_rect(slide, Inches(6.8), iy, Inches(0.05), Inches(0.72), fill=C["yellow"])
        add_text(slide, s, Inches(7.0), iy + Inches(0.1), Inches(5.5), Inches(0.55),
                 font_size=Pt(9.5), color=C["white"], wrap=True)

    # Next steps bar
    ns_top = Inches(5.6)
    add_rect(slide, Inches(0.5), ns_top, Inches(12.33), Inches(1.5),
             fill=C["surface"], line=C["surface2"])
    add_text(slide, "RECOMMENDED NEXT STEPS", Inches(0.65), ns_top + Inches(0.1),
             Inches(5), Inches(0.2), font_size=Pt(8), bold=True, color=C["muted"])
    steps = [
        "› Monitor P99 spikes monthly to catch regressions early",
        "› Set alert thresholds on backend 500 error rate",
        "› Correlate high-latency months with deployment history",
        "› Track hit growth to plan capacity for next quarter",
    ]
    for i, s in enumerate(steps):
        col = i % 2
        row = i // 2
        sx = Inches(0.65 + col * 6.2)
        sy = ns_top + Inches(0.35 + row * 0.45)
        add_text(slide, s, sx, sy, Inches(6.0), Inches(0.35),
                 font_size=Pt(9), color=C["muted"])


# ── Main ──────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    fe_keys = [k for k, v in DATA.items() if v["type"] == "frontend"]
    be_keys = [k for k, v in DATA.items() if v["type"] == "backend"]
    total = 2 + len(fe_keys) + len(be_keys) + 2

    print(f"Building slide 1/{total}: Title...")
    build_title_slide(prs)

    print(f"Building slide 2/{total}: Executive Summary...")
    build_summary_slide(prs)

    for i, k in enumerate(fe_keys):
        print(f"Building slide {3+i}/{total}: {DATA[k]['label']}...")
        build_function_slide(prs, k, "Frontend Performance", is_backend=False)

    for i, k in enumerate(be_keys):
        print(f"Building slide {3+len(fe_keys)+i}/{total}: {DATA[k]['label']}...")
        build_function_slide(prs, k, "Backend Performance", is_backend=True)

    print(f"Building slide {total-1}/{total}: Comparison...")
    build_comparison_slide(prs)

    print(f"Building slide {total}/{total}: Takeaways...")
    build_takeaways_slide(prs)

    out = "cexp_presentation.pptx"
    prs.save(out)
    print(f"\n✓  Saved: {out}")
    print("  Open it in PowerPoint or Keynote.")

if __name__ == "__main__":
    main()
